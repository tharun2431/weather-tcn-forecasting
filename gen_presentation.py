"""Generate Tharun_Milestone1.pptx — Milestone 1 Presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colour palette ──
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)    # Deep navy
BG_CARD   = RGBColor(0x16, 0x21, 0x3E)    # Card navy
ACCENT    = RGBColor(0x00, 0xD2, 0xFF)    # Cyan accent
ACCENT2   = RGBColor(0x7C, 0x3A, 0xED)    # Purple accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xB0, 0xBC, 0xD4)    # Muted text
GREEN     = RGBColor(0x10, 0xB9, 0x81)
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return tf

def add_accent_line(slide, left, top, width, color=ACCENT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_section_header(slide, number, title):
    add_text(slide, 0.8, 0.4, 1.5, 0.6, f"0{number}", font_size=48,
             color=ACCENT, bold=True)
    add_text(slide, 2.2, 0.5, 9, 0.5, title, font_size=32,
             color=WHITE, bold=True)
    add_accent_line(slide, 0.8, 1.1, 11)

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

# Decorative top line
add_accent_line(slide, 0, 0, 13.333, ACCENT)

# Module info
add_text(slide, 0.8, 1.0, 12, 0.4,
         "MSc Data Science  |  Deep Learning Applications  |  CMP-L016",
         font_size=14, color=LIGHT)

# Title
add_text(slide, 0.8, 2.0, 12, 1.2,
         "Hybrid Deep Learning Models for\nWeather Prediction",
         font_size=44, color=WHITE, bold=True)

# Subtitle line
add_accent_line(slide, 0.8, 3.6, 5, ACCENT)

add_text(slide, 0.8, 3.9, 6, 0.4,
         "Project #28 — TCN, Hybrid & Ensemble Models",
         font_size=20, color=ACCENT)

# Author info card
add_card(slide, 0.8, 5.0, 5, 1.5)
add_text(slide, 1.1, 5.1, 4.5, 0.4,
         "Tharun", font_size=24, color=WHITE, bold=True)
add_text(slide, 1.1, 5.6, 4.5, 0.3,
         "Milestone 1 Presentation", font_size=14, color=LIGHT)
add_text(slide, 1.1, 6.0, 4.5, 0.3,
         "February 2026", font_size=14, color=LIGHT)

# Right side decorative element
add_card(slide, 8.5, 4.5, 4.2, 2.5, RGBColor(0x10, 0x1D, 0x37))
add_text(slide, 8.8, 4.7, 3.8, 0.4,
         "📊 Jena Climate Dataset", font_size=16, color=ACCENT)
add_text(slide, 8.8, 5.2, 3.8, 0.3,
         "760,000+ observations", font_size=13, color=LIGHT)
add_text(slide, 8.8, 5.6, 3.8, 0.3,
         "14 meteorological features", font_size=13, color=LIGHT)
add_text(slide, 8.8, 6.0, 3.8, 0.3,
         "7 years (2009–2022)", font_size=13, color=LIGHT)
add_text(slide, 8.8, 6.4, 3.8, 0.3,
         "10-minute recording interval", font_size=13, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda / Overview
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 0, "Presentation Overview")

items = [
    ("01", "Background & Problem Definition", "Why weather prediction matters"),
    ("02", "Problem Statement", "Defining our deep learning challenge"),
    ("03", "Proposed Solution: Hybrid DL", "TCN-LSTM hybrid + Stacking Ensemble"),
    ("04", "Model Architecture Details", "4 models: 2 base + 1 hybrid + 1 ensemble"),
    ("05", "Dataset: Jena Climate", "Data selection and characteristics"),
    ("06", "Methodology & Pipeline", "End-to-end workflow"),
    ("07", "Expected Outcomes", "What we aim to demonstrate"),
    ("08", "References", "Supporting literature"),
]

for i, (num, title, desc) in enumerate(items):
    row = i // 2
    col = i % 2
    x = 1.0 + col * 6.0
    y = 1.8 + row * 1.3

    add_card(slide, x, y, 5.5, 1.1)
    add_text(slide, x + 0.2, y + 0.15, 0.6, 0.4, num,
             font_size=22, color=ACCENT, bold=True)
    add_text(slide, x + 0.9, y + 0.1, 4.2, 0.4, title,
             font_size=16, color=WHITE, bold=True)
    add_text(slide, x + 0.9, y + 0.55, 4.2, 0.4, desc,
             font_size=12, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: Background & Context
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 1, "Background & Context")

# Left card
add_card(slide, 0.8, 1.6, 5.5, 5.2)
add_text(slide, 1.1, 1.8, 5, 0.4,
         "Why Weather Prediction Matters", font_size=20, color=ACCENT, bold=True)
add_bullet_list(slide, 1.1, 2.4, 5, 3.5, [
    "▸  Weather affects agriculture, energy, transport,\n    and disaster management globally",
    "▸  Accurate short-term forecasting saves lives\n    and billions in economic value annually",
    "▸  Climate change increases weather volatility,\n    making prediction more critical than ever",
    "▸  Traditional Numerical Weather Prediction (NWP)\n    requires massive computational resources",
], font_size=14)

# Right card
add_card(slide, 7.0, 1.6, 5.5, 5.2)
add_text(slide, 7.3, 1.8, 5, 0.4,
         "The Deep Learning Opportunity", font_size=20, color=GREEN, bold=True)
add_bullet_list(slide, 7.3, 2.4, 5, 3.5, [
    "▸  Deep learning models learn directly from\n    historical observation data",
    "▸  No need for complex physics simulations\n    — purely data-driven approach",
    "▸  Real-time inference: predictions in milliseconds\n    vs hours for NWP models",
    "▸  Can capture non-linear relationships between\n    atmospheric variables automatically",
], font_size=14)

add_text(slide, 0.8, 6.5, 12, 0.4,
         "\"Deep learning is transforming weather forecasting from physics simulations to data-driven predictions\"",
         font_size=13, color=LIGHT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: Problem Definition
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 2, "Problem Definition")

# Main problem statement card
add_card(slide, 0.8, 1.6, 11.7, 1.4, RGBColor(0x1A, 0x1A, 0x4E))
add_text(slide, 1.2, 1.7, 11, 0.5,
         "Research Question", font_size=14, color=ACCENT)
add_text(slide, 1.2, 2.1, 11, 0.6,
         "Can hybrid and ensemble deep learning architectures outperform standalone models\n"
         "for multivariate weather time-series forecasting?",
         font_size=18, color=WHITE, bold=True)

# Three sub-cards
cards = [
    ("Input", "7 days (168 hours) of\n14 atmospheric variables\nrecorded hourly", ACCENT),
    ("Output", "Temperature prediction\n1–24 hours into the future\nin degrees Celsius", GREEN),
    ("Goal", "Prove hybrid TCN-LSTM\nvs LSTM and TCN baselines\nvia ensemble combination", ORANGE),
]

for i, (title, desc, color) in enumerate(cards):
    x = 0.8 + i * 4.1
    add_card(slide, x, 3.5, 3.8, 2.5)
    add_text(slide, x + 0.3, 3.7, 3.2, 0.4, title,
             font_size=20, color=color, bold=True)
    add_accent_line(slide, x + 0.3, 4.2, 2.0, color)
    add_text(slide, x + 0.3, 4.4, 3.2, 1.5, desc,
             font_size=14, color=LIGHT)

# Why DL is appropriate
add_card(slide, 0.8, 6.3, 11.7, 0.9, RGBColor(0x0E, 0x2A, 0x1E))
add_text(slide, 1.2, 6.4, 11, 0.7,
         "Why Deep Learning?  Weather data is sequential, multivariate, and non-linear — "
         "ideal for deep sequence models. TCNs offer parallelisable training with large "
         "receptive fields, overcoming RNN limitations.",
         font_size=13, color=GREEN)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: Proposed Solution — TCN Architecture
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 3, "Proposed Solution — Hybrid Architecture")

# TCN description card
add_card(slide, 0.8, 1.6, 6.5, 5.2)
add_text(slide, 1.1, 1.8, 6, 0.4,
         "Temporal Convolutional Network (TCN)", font_size=18, color=ACCENT, bold=True)

add_bullet_list(slide, 1.1, 2.4, 6, 4.0, [
    "▸  1D Causal Convolutions\n    Ensures no information leakage from future time steps",
    "▸  Dilated Convolutions (d = 1, 2, 4, 8, 16)\n    Exponentially growing receptive field without\n    increasing parameters — covers 125+ hours",
    "▸  Residual Connections\n    Stable gradient flow through deep networks,\n    enables training of 5+ layer architectures",
    "▸  Weight Normalisation + BatchNorm\n    Faster convergence and training stability",
], font_size=13)

# Architecture specs card
add_card(slide, 7.8, 1.6, 4.7, 2.3)
add_text(slide, 8.1, 1.8, 4.2, 0.4,
         "Architecture Specs", font_size=16, color=GREEN, bold=True)

specs = [
    ("Layers", "5 residual blocks"),
    ("Filters", "64 per layer"),
    ("Kernel", "3 (causal)"),
    ("Dilations", "1, 2, 4, 8, 16"),
    ("Receptive Field", "125 hours (5.2 days)"),
    ("Parameters", "~85K"),
]
for i, (k, v) in enumerate(specs):
    y = 2.25 + i * 0.28
    add_text(slide, 8.1, y, 2.0, 0.3, k, font_size=11, color=LIGHT)
    add_text(slide, 10.0, y, 2.5, 0.3, v, font_size=11, color=WHITE, bold=True)

# Unique aspects card
add_card(slide, 7.8, 4.2, 4.7, 2.6, RGBColor(0x1A, 0x1A, 0x4E))
add_text(slide, 8.1, 4.4, 4.2, 0.4,
         "★ Unique Aspects", font_size=16, color=ORANGE, bold=True)
add_bullet_list(slide, 8.1, 4.9, 4.2, 2.0, [
    "✦  Parallelisable — unlike sequential\n    LSTM/GRU processing",
    "✦  No vanishing gradient problem",
    "✦  TCN-LSTM hybrid combines local feature\n    extraction with sequential memory",
    "✦  Stacking ensemble: meta-learner\n    optimally combines all models",
], font_size=12, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: Model Comparison Strategy
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 4, "4-Model Comparison Strategy")

models = [
    ("LSTM", "Recurrent Baseline", "2-layer LSTM (128 units) → Dense(64) → 1\nGated memory for long-range dependencies",
     RGBColor(0x34, 0x98, 0xDB), "~200K params"),
    ("TCN", "Convolutional Baseline", "5 residual blocks, dilated causal conv\nParallelisable, 125h receptive field",
     ACCENT, "~85K params"),
    ("TCN-LSTM", "★ Hybrid (Key Innovation)", "TCN encoder → LSTM decoder\nLocal features + sequential memory",
     ORANGE, "~120K params"),
    ("Ensemble", "★ Stacking", "Ridge meta-learner on all model\npredictions — optimal combination",
     ACCENT2, "Meta-learner"),
]

for i, (name, subtitle, desc, color, params) in enumerate(models):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.2
    y = 1.6 + row * 2.8
    add_card(slide, x, y, 5.8, 2.5)
    add_text(slide, x + 0.3, y + 0.2, 5.0, 0.5, name,
             font_size=28, color=color, bold=True)
    add_text(slide, x + 0.3, y + 0.7, 5.0, 0.3, subtitle,
             font_size=12, color=LIGHT)
    add_accent_line(slide, x + 0.3, y + 1.05, 3.5, color)
    add_text(slide, x + 0.3, y + 1.2, 5.0, 0.8, desc,
             font_size=13, color=LIGHT)
    add_text(slide, x + 0.3, y + 2.1, 5.0, 0.3, params,
             font_size=12, color=color, bold=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: Dataset — Jena Climate
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 5, "Dataset: Jena Climate 2009–2022")

# Dataset overview card
add_card(slide, 0.8, 1.6, 5.5, 2.5)
add_text(slide, 1.1, 1.8, 5, 0.4,
         "Dataset Overview", font_size=18, color=ACCENT, bold=True)

details = [
    ("Source", "Max Planck Institute for Biogeochemistry"),
    ("Location", "Jena, Germany"),
    ("Period", "January 2009 – December 2016"),
    ("Raw Rows", "420,551 (every 10 minutes)"),
    ("Resampled", "~70,000 (hourly)"),
    ("Features", "14 meteorological variables"),
    ("Target", "T (degC) — Temperature"),
]
for i, (k, v) in enumerate(details):
    y = 2.3 + i * 0.25
    add_text(slide, 1.1, y, 2.2, 0.25, k, font_size=11, color=LIGHT)
    add_text(slide, 3.2, y, 3.0, 0.25, v, font_size=11, color=WHITE, bold=True)

# Features card
add_card(slide, 7.0, 1.6, 5.5, 2.5)
add_text(slide, 7.3, 1.8, 5, 0.4,
         "14 Meteorological Features", font_size=18, color=GREEN, bold=True)

features = [
    "Temperature (°C), Potential temp (K), Dew point (°C)",
    "Atmospheric pressure (mbar)",
    "Relative humidity (%)",
    "Vapour pressure: max, actual, deficit (mbar)",
    "Specific humidity (g/kg), H₂O concentration",
    "Air density (g/m³)",
    "Wind speed, max wind speed (m/s), direction (°)",
]
add_bullet_list(slide, 7.3, 2.4, 5, 2.0, [f"▸  {f}" for f in features], font_size=11)

# Suitability card
add_card(slide, 0.8, 4.4, 5.5, 2.8)
add_text(slide, 1.1, 4.6, 5, 0.4,
         "Why This Dataset?", font_size=18, color=ORANGE, bold=True)
add_bullet_list(slide, 1.1, 5.1, 5, 2.0, [
    "▸  Publicly available and well-documented",
    "▸  Used in TensorFlow's official time-series tutorial",
    "▸  Rich multivariate structure (14 inter-related features)",
    "▸  Strong seasonal & diurnal patterns — ideal for\n    demonstrating sequence model capabilities",
    "▸  7 years provides sufficient data for robust\n    train/validation/test splitting",
], font_size=12)

# Challenges card
add_card(slide, 7.0, 4.4, 5.5, 2.8, RGBColor(0x2A, 0x1A, 0x1A))
add_text(slide, 7.3, 4.6, 5, 0.4,
         "Data Challenges", font_size=18, color=RGBColor(0xFF, 0x63, 0x63), bold=True)
add_bullet_list(slide, 7.3, 5.1, 5, 2.0, [
    "▸  Erroneous values: negative wind speeds require\n    cleaning (set to 0)",
    "▸  High volume: 420K rows needs resampling to\n    hourly for computational feasibility",
    "▸  Non-stationarity: seasonal trends must be\n    handled through normalisation",
    "▸  Temporal ordering: must use chronological\n    splitting (no shuffling) to prevent data leakage",
], font_size=12)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 8: Methodology & Pipeline
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 6, "Methodology & Training Pipeline")

# Pipeline steps
steps = [
    ("1", "Data\nPreprocessing", "Clean, resample\nto hourly, fix\nerroneous values", ACCENT),
    ("2", "Feature\nEngineering", "StandardScaler\n(train-only fit),\n14 features", GREEN),
    ("3", "Window\nCreation", "168h sliding window\n→ 1h prediction\n(DataLoaders)", ORANGE),
    ("4", "Model\nTraining", "4 architectures:\n2 base + 1 hybrid\n+ 1 ensemble", RGBColor(0xE7, 0x4C, 0x3C)),
    ("5", "HP Tuning\n& Ablation", "27-config grid\nfor TCN, depth\nablation study", ACCENT2),
    ("6", "Evaluation\n& Analysis", "Metrics, residuals\nmulti-step forecast\ncomparison", GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.1
    add_card(slide, x, 1.8, 1.9, 3.5)
    add_text(slide, x + 0.15, 1.9, 1.6, 0.45, num,
             font_size=32, color=color, bold=True)
    add_text(slide, x + 0.15, 2.35, 1.6, 0.7, title,
             font_size=13, color=WHITE, bold=True)
    add_accent_line(slide, x + 0.15, 3.1, 1.3, color)
    add_text(slide, x + 0.15, 3.3, 1.6, 1.5, desc,
             font_size=10, color=LIGHT)

# Bottom: Data split info
add_card(slide, 0.8, 5.6, 11.7, 1.5)
add_text(slide, 1.1, 5.7, 3, 0.4,
         "Chronological Data Split", font_size=16, color=ACCENT, bold=True)

splits = [
    ("Train (70%)", "2009–2014", "~49K hours", GREEN),
    ("Validation (15%)", "2014–2015", "~10.5K hours", ORANGE),
    ("Test (15%)", "2015–2016", "~10.5K hours", RGBColor(0xE7, 0x4C, 0x3C)),
]
for i, (name, period, size, color) in enumerate(splits):
    x = 1.1 + i * 3.8
    add_text(slide, x, 6.15, 3.5, 0.3, name, font_size=13, color=color, bold=True)
    add_text(slide, x, 6.45, 3.5, 0.25, f"{period}  •  {size}", font_size=11, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 9: Expected Outcomes
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 7, "Expected Outcomes")

outcomes = [
    ("Quantitative Comparison", 
     "MAE, RMSE, R², MAPE — all models\nevaluated on identical test set",
     ACCENT),
    ("Training Analysis",
     "Learning curves, convergence rates,\ntraining time comparison",
     GREEN),
    ("Ablation Study",
     "Impact of TCN depth (2–5 blocks)\non accuracy and receptive field",
     ORANGE),
    ("Multi-step Forecasting",
     "24h autoregressive prediction,\nerror accumulation analysis",
     ACCENT2),
    ("Residual Analysis",
     "Error distributions, seasonal and\ndiurnal error patterns",
     RGBColor(0xE7, 0x4C, 0x3C)),
    ("Live Deployment",
     "Web application for real-time\nweather prediction using\nhybrid ensemble model",
     RGBColor(0xF5, 0x9E, 0x0B)),
]

for i, (title, desc, color) in enumerate(outcomes):
    row = i // 3
    col = i % 3
    x = 0.8 + col * 4.1
    y = 1.6 + row * 2.8

    add_card(slide, x, y, 3.8, 2.3)
    add_text(slide, x + 0.3, y + 0.2, 3.3, 0.4, title,
             font_size=16, color=color, bold=True)
    add_accent_line(slide, x + 0.3, y + 0.7, 2.5, color)
    add_text(slide, x + 0.3, y + 0.9, 3.3, 1.2, desc,
             font_size=13, color=LIGHT)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10: References
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_section_header(slide, 8, "References")

refs = [
    "[1]  S. Bai, J.Z. Kolter, V. Koltun, \"An Empirical Evaluation of Generic Convolutional and\n"
    "      Recurrent Networks for Sequence Modeling,\" arXiv:1803.01271, 2018.",
    
    "[2]  S. Hochreiter and J. Schmidhuber, \"Long Short-Term Memory,\" Neural Computation,\n"
    "      vol. 9, no. 8, pp. 1735–1780, 1997.",
    
    "[3]  K. Cho et al., \"Learning Phrase Representations using RNN Encoder-Decoder for Statistical\n"
    "      Machine Translation,\" arXiv:1406.1078, 2014.",
    
    "[4]  K. He, X. Zhang, S. Ren, and J. Sun, \"Deep Residual Learning for Image Recognition,\"\n"
    "      Proc. CVPR, pp. 770–778, 2016.",
    
    "[5]  A. van den Oord et al., \"WaveNet: A Generative Model for Raw Audio,\"\n"
    "      arXiv:1609.03499, 2016.",
    
    "[6]  A. Paszke et al., \"PyTorch: An Imperative Style, High-Performance Deep Learning Library,\"\n"
    "      NeurIPS, 2019.",
    
    "[7]  Max Planck Institute for Biogeochemistry, \"Jena Climate Dataset 2009–2022,\"\n"
    "      Available: https://www.bgc-jena.mpg.de/wetter/",
]

add_card(slide, 0.8, 1.5, 11.7, 5.5)
for i, ref in enumerate(refs):
    y = 1.7 + i * 0.72
    add_text(slide, 1.2, y, 11, 0.65, ref,
             font_size=11, color=LIGHT)

add_text(slide, 0.8, 6.8, 12, 0.4,
         "All references follow IEEE citation format as required",
         font_size=11, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11: Thank You
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_accent_line(slide, 0, 0, 13.333, ACCENT)

add_text(slide, 1.0, 2.5, 11.3, 1.0,
         "Thank You", font_size=52, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(slide, 1.0, 3.5, 11.3, 0.5,
         "Questions & Discussion", font_size=24, color=ACCENT,
         alignment=PP_ALIGN.CENTER)

add_accent_line(slide, 4.5, 4.2, 4.3, ACCENT)

add_card(slide, 3.5, 4.8, 6.3, 1.8, BG_CARD)
add_text(slide, 3.8, 5.0, 5.7, 0.4,
         "Tharun", font_size=20, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, 3.8, 5.5, 5.7, 0.3,
         "MSc Data Science — Deep Learning Applications", font_size=14, color=LIGHT,
         alignment=PP_ALIGN.CENTER)
add_text(slide, 3.8, 5.9, 5.7, 0.3,
         "Project #28 — Weather Prediction with Hybrid DL", font_size=14, color=ACCENT,
         alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = os.path.join("outputs", "Tharun_Milestone1.pptx")
os.makedirs("outputs", exist_ok=True)
prs.save(output_path)
print(f"✓ Presentation saved: {output_path}")
print(f"  Slides: {len(prs.slides)}")

# Also save to Desktop for easy access
desktop_path = os.path.join(os.path.expanduser("~"), "Desktop", "Tharun_Milestone1.pptx")
prs.save(desktop_path)
print(f"✓ Also saved to Desktop: {desktop_path}")
